import fitz  # PyMuPDF
import pytesseract
import easyocr
from PIL import Image
import os
from typing import List, Dict, Any, Optional
from pathlib import Path

# Imports for semantic chunking
from langchain_experimental.text_splitter import SemanticChunker
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.document_loaders import CSVLoader, Docx2txtLoader, WebBaseLoader, OnlinePDFLoader
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd
import json
import zipfile
import xml.etree.ElementTree as ET
from docx import Document
import requests
import re
import base64
from langchain_groq import ChatGroq
from langchain_ollama import ChatOllama
from io import BytesIO

import logging

# Local application imports
# from backend.app.core.config import settings  #NO NEED 
from backend.app import OLLAMA_MODEL, OLLAMA_BASE_URL, GROQ_API_KEY, GROQ_MODEL,EMBEDDING_MODEL_NAME, OLLAMA_CHAT_MODEL
from backend.app.services.vstore_svc import VectorStoreService
from backend.app.services.system_message import QUERY,TITLE_PROMPT
from backend.app.types.response_format import TitleResponse,VLMResponse

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


class DocParserFastService:
    """
    Service to parse documents (PDFs, images), extract text,
    perform structure-aware semantic chunking, and add chunks to the vector store.
    
    Refactored for performance, correctness, and robustness.
    """
    def __init__(self, vector_store_service: VectorStoreService):
        self.vector_store_service = vector_store_service
        
        # --- PERFORMANCE: Initialize models only ONCE ---
        self.text_llm = ChatGroq(
            model = GROQ_MODEL,
            api_key = GROQ_API_KEY,
        )
        self.title_chain = TITLE_PROMPT | self.text_llm.with_structured_output(TitleResponse)
        self.vlm = None
        # self._initialize_models()

        # Initialize embeddings and splitter
        try:
            self.embeddings = HuggingFaceEmbeddings(
                model_name=EMBEDDING_MODEL_NAME,
                model_kwargs={'device': 'cpu'},
                encode_kwargs={"normalize_embeddings": True}
            )
            
            self.semantic_splitter = SemanticChunker(
                    self.embeddings,
                    breakpoint_threshold_type="percentile",  # or "standard_deviation"
                    breakpoint_threshold_amount=95,          # controls sensitivity
                    min_chunk_size=1000,                      # in characters
            )
            logger.info(f"Successfully initialized embedding model: {EMBEDDING_MODEL_NAME} and Semantic Chunker")
        except Exception as e:
            logger.error(f"Error initializing embedding model: {e}. Chunking may be suboptimal.")
            self.embeddings = None
            self.semantic_splitter = RecursiveCharacterTextSplitter(
                chunk_size=800,
                chunk_overlap=200,
                length_function=len,
                is_separator_regex=False,
                separators=["\n\n", "\n", ". ", "? ", "! ", ",", " ", ""]
            )

    def _initialize_models(self):
        """Initializes LLM and VLM models to be reused across the service."""
        try:
            # For titles and other text tasks
            self.text_llm = ChatOllama(model=OLLAMA_CHAT_MODEL, base_url=OLLAMA_BASE_URL)
            self.title_chain = TITLE_PROMPT | self.text_llm.with_structured_output(TitleResponse)

            # For Vision-Language Model tasks (e.g., image analysis)
            self.vlm = ChatOllama(model=OLLAMA_MODEL, base_url=OLLAMA_BASE_URL)
            self.vlm_structured = self.vlm.with_structured_output(VLMResponse)
            logger.info("LLM and VLM models initialized successfully.")
        except Exception as e:
            logger.error(f"Failed to initialize LLM/VLM models: {e}")

    def _get_title_from_llm(self, paragraphs: str) -> Optional[TitleResponse]:
        if not self.title_chain:
            logger.error("Title generation chain not initialized.")
            return None
        try:
            result = self.title_chain.invoke({"paragraphs": paragraphs})
            return result.title
        except Exception as e:
            logger.error(f"Error getting title from LLM: {e}")
            return None

    def _extract_text_from_vlm(self, encoded_image):
        if not self.vlm_structured:
            logger.error("VLM model not initialized.")
            return None
        try:
            # Decode and convert to PIL image
            decoded_bytes = base64.b64decode(encoded_image.encode("utf-8"))
            img = Image.open(BytesIO(decoded_bytes))

            # OCR
            reader = easyocr.Reader(['en'], gpu=False)
            text_from_ocr = reader.readtext(decoded_bytes, detail=0)
            # text_from_ocr = pytesseract.image_to_string(img)

            # Generate query using OCR text
            generated_query = QUERY.invoke({"context": text_from_ocr})

            # Build prompt
            prompt = [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": generated_query
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{encoded_image}",
                            },
                        },
                    ],
                }
            ]

            # Call VLM
            content = self.vlm_structured.invoke(prompt)
            logging.info("VLM and Pytesseract OCR extracted data successfully.")
            logging.info(f"VLM Data:\n{content}")
            return content

        except Exception as e:
            # Fallback OCR only
            text_from_ocr : str
            try:
                decoded_bytes = base64.b64decode(encoded_image.encode("utf-8"))
                img = Image.open(BytesIO(decoded_bytes))
                text_from_ocr = pytesseract.image_to_string(img)
            except Exception as inner_e:
                logger.error(f"Error decoding image in fallback: {inner_e}")
                return {"title": "Error", "description": "Failed to decode image."}

            title = self._get_title_from_llm(text_from_ocr)
            logger.error(f"Error extracting text with VLM: {e}")
            logging.info("Fallback: using Tesseract OCR only.")

            return {
                "title": title,  # typo fixed
                "description": f"{title}:\n{text_from_ocr}"
            }


    def _is_heading(self, text: str) -> bool:
        """
        Detect if a line/paragraph is a heading based on heuristic:
        - Mostly uppercase or capitalized and short (< 8 words)
        """
        words = text.strip().split()
        if len(words) <= 8 and sum(1 for w in words if w.isupper()) >= len(words) * 0.6:
            return True
        return False

    # def _extract_sections_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
    #     """
    #     Extracts structured sections from a PDF using text and image understanding.
    #     Heuristics:
    #     - Uses _is_heading() to detect section titles.
    #     - If both image and text are present, render the entire page as an image and analyze it using VLM.
    #     - If only text, use text+LLM title.
    #     - If only image(s), use individual image analysis via VLM.
    #     - Avoids duplicate sections.
    #     """
    #     sections = []
    #     doc = fitz.open(file_path)
    #     section_counter = 0
    #     seen_contents = set()

    #     for page_index in range(len(doc)):
    #         page = doc.load_page(page_index)
    #         blocks = page.get_text("blocks", sort=True)

    #         # Step 1: Parse text blocks and headings
    #         text_paragraphs = []
    #         heading_text = None
    #         for block in blocks:
    #             if block[6] != 0:
    #                 continue
    #             text = block[4].replace('\r', ' ').replace('\n', ' ').strip()
    #             if not text:
    #                 continue
    #             if self._is_heading(text) and not heading_text:
    #                 heading_text = text
    #             else:
    #                 text_paragraphs.append(text)

    #         has_text = bool(text_paragraphs)
    #         has_heading = heading_text is not None
    #         images = page.get_images(full=True)
    #         has_images = bool(images)

    #         def add_section(title: str, paragraphs: list):
    #             nonlocal section_counter
    #             content_key = title + " " + " ".join(paragraphs)
    #             if content_key not in seen_contents and paragraphs:
    #                 section_counter += 1
    #                 sections.append({
    #                     "title": title,
    #                     "page_number": page_index + 1,
    #                     "paragraphs": paragraphs,
    #                     "section_index": section_counter
    #                 })
    #                 seen_contents.add(content_key)

    #         # CASE 1: Chart-style full-page rendering (text + images present)
    #         if has_text and has_images:
    #             pix = page.get_pixmap(dpi=200)
    #             img_bytes = pix.tobytes("png")
    #             img_base64 = base64.b64encode(img_bytes).decode("utf-8")
    #             content = self._extract_text_from_vlm(encoded_image=img_base64)
    #             if content:
    #                 add_section(content.title or "Chart Section", [content.description])

    #         # CASE 2: Text-only pages
    #         elif has_text:
    #             title = heading_text or ""
    #             if not title:
    #                 llm_response = self._get_title_from_llm(" ".join(text_paragraphs))
    #                 title = llm_response.title if llm_response else "Text Section"
    #             add_section(title, text_paragraphs)

    #         # CASE 3: Image-only pages — handle each image separately
    #         elif has_images:
    #             for img_info in images:
    #                 xref = img_info[0]
    #                 pix = fitz.Pixmap(doc, xref)
    #                 if pix.n > 4:
    #                     pix = fitz.Pixmap(fitz.csRGB, pix)
    #                 img_bytes = pix.tobytes("png")
    #                 img_base64 = base64.b64encode(img_bytes).decode("utf-8")

    #                 content = self._extract_text_from_vlm(encoded_image=img_base64)
    #                 if content:
    #                     add_section(content.title or "Image Section", [content.description])

    #     doc.close()
    #     logger.info("Text is extracted from PDF sucessfully.....")
    #     return sections
    
    def _extract_sections_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts structured sections from a PDF by rendering each page as an image and analyzing it using VLM.
        Ignores text blocks entirely and uses Vision-Language Model (VLM) on full-page image.
        """
        sections = []
        doc = fitz.open(file_path)
        section_counter = 0
        seen_contents = set()

        for page_index in range(len(doc)):
            page = doc.load_page(page_index)

            # Render full page as image
            pix = page.get_pixmap(dpi=200) 
            img_bytes = pix.tobytes("png")
            img_base64 = base64.b64encode(img_bytes).decode("utf-8")

            # Extract text using Vision-Language Model
            content = self._extract_text_from_vlm(encoded_image=img_base64)

            def add_section(title: str, paragraphs: list):
                nonlocal section_counter
                content_key = title + " " + " ".join(paragraphs)
                if content_key not in seen_contents and paragraphs:
                    section_counter += 1
                    sections.append({
                        "title": title,
                        "page_number": page_index + 1,
                        "paragraphs": paragraphs,
                        "section_index": section_counter
                    })
                    seen_contents.add(content_key)

            if content:
                title = content.title
                description = content.description.strip()
                if description:
                    add_section(title, [description])
                else:
                    logger.warning(f"Empty description from VLM on page {page_index + 1}")
            else:
                logger.warning(f"VLM failed to extract content from page {page_index + 1}")

        doc.close()
        logger.info("Text is extracted from PDF via VLM rendering successfully.")
        return sections

    def _extract_text_from_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts all content from each PPTX slide as one section:
        - Uses _is_heading to set section title.
        - Combines text, chart summaries, and image descriptions into one 'paragraphs' list.
        """
        prs = Presentation(file_path)
        sections = []
        section_counter = 0

        for i, slide in enumerate(prs.slides):
            section_counter += 1
            paragraphs = []

            section_title = f"slide_{i+1}"
            for shape in slide.shapes:
                # 1. Text content
                if hasattr(shape, "text"):
                    text = shape.text.strip().replace('\r', ' ').replace('\n', ' ')
                    if not text:
                        continue
                    if self._is_heading(text) and section_title.startswith("slide_"):
                        section_title = text  # First heading becomes title
                    else:
                        paragraphs.append(text)

                # 2. Chart content
                elif isinstance(shape, GraphicFrame) and shape.has_chart:
                    chart = shape.chart
                    for series in chart.series:
                        values = [v for v in series.values]
                        chart_text = (
                            f"Chart: Series '{series.name}', values: {values}"
                        )
                        paragraphs.append(chart_text)

                # 3. Image content (VLM-based)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = shape.image
                        img_bytes = img.blob
                        img_base64 = base64.b64encode(img_bytes).decode("utf-8")
                        content = self._extract_text_from_vlm(encoded_image=img_base64)
                        if content:
                            paragraphs.append(f"Image Insight: {content.description}")
                            # Optionally override title if VLM title is better
                            if len(content.title or "") > len(section_title):
                                section_title = content.title
                            
                    except Exception as e:
                        print(f"unable to process the document---Error : {e}")
                        continue

            if paragraphs:
                try:
                    llm_title = self._get_title_from_llm("\n\n".join(paragraphs))
                    if llm_title and len(llm_title.strip()) > 0:
                        section_title = llm_title
                except Exception as e:
                    print(f"Error getting title from LLM for slide {i + 1}: {e}")
                    # Keep the existing section_title
                
                sections.append({
                    "title": section_title,
                    "page_number": i + 1,
                    "paragraphs": paragraphs,
                    "section_index": section_counter
                })

        return sections


    def _extract_text_from_csv(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from a CSV file and groups rows into sections.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0

        loader = CSVLoader(file_path=file_path)
        # documents = loader.load_and_split(text_splitter=self.semantic_splitter)
        documents = loader.load()

        for i, doc in enumerate(documents):
            text = doc.page_content.strip()
            title = self._get_title_from_llm("CSV DATA : "+ text)
            if not text:
                continue
            current_section = {
                "title": title,
                "page_number": i + 1,
                "paragraphs": [text]
            }
            section_counter += 1
            current_section["section_index"] = section_counter
            sections.append(current_section)

        return sections

    def _extract_text_from_excel(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from Excel file and groups each sheet as a section.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0

        excel_file = pd.ExcelFile(file_path)

        for i, sheet_name in enumerate(excel_file.sheet_names):
            df = pd.read_excel(file_path, sheet_name=sheet_name)

            # Convert DataFrame to text (one big paragraph)
            text_content = df.to_string(index=False).strip()
            if not text_content:
                continue

            current_section = {
                "title": sheet_name if sheet_name else f"sheet_{i+1}_untitled",
                "page_number": i + 1,
                "paragraphs": [text_content],
            }

            section_counter += 1
            current_section["section_index"] = section_counter
            sections.append(current_section)

        return sections
    
    def _extract_text_from_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from DOCX file and split into sections.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0
        loader = Docx2txtLoader(file_path=file_path)
        documents = loader.load()
        
        for i, doc in enumerate(documents):
            text = doc.page_content.strip()
            if not text:
                continue
            section_counter += 1
            sections.append({
                "title": f"page_{i+1}_untitled",
                "page_number": i + 1,
                "paragraphs": [text],
                "section_index": section_counter
            })

        return sections

    def _extract_text_from_json(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Handles JSON files containing multiple JSON objects (one per line or concatenated).
        Extracts structured sections with 'title', 'page_number', 'paragraphs', and 'section_index'.
        """
        sections = []
        section_counter = 0

        with open(file_path, "r", encoding="utf-8") as f:
            for i, line in enumerate(f):
                line = line.strip()
                if not line:
                    continue
                try:
                    item = json.loads(line)
                except json.JSONDecodeError as e:
                    print(f"Skipping line {i+1} due to JSON decode error: {e}")
                    continue

                title = item.get("title") or f"section_{i+1}_untitled"
                text_parts = []

                for key, value in item.items():
                    if isinstance(value, str):
                        text_parts.append(value)
                    elif isinstance(value, list):
                        text_parts.extend([str(v) for v in value if isinstance(v, str)])

                paragraph_text = " ".join(text_parts).strip()
                if not paragraph_text:
                    continue

                section_counter += 1
                current_section = {
                    "title": title,
                    "page_number": section_counter,
                    "paragraphs": [paragraph_text],
                    "section_index": section_counter
                }
                sections.append(current_section)

        return sections


    def _extract_text_from_txt(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from TXT file and splits by headings or paragraphs.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0

        with open(file_path, "r", encoding="utf-8") as f:
            lines = [line.strip() for line in f if line.strip()]

        current_section = {
            "title": "page_1_untitled",
            "page_number": 1,
            "paragraphs": []
        }

        for line in lines:
            if self._is_heading(line):
                if current_section["paragraphs"]:
                    sections.append(current_section)
                section_counter += 1
                current_section = {
                    "title": line,
                    "page_number": 1,
                    "paragraphs": [],
                    "section_index": section_counter
                }
            else:
                current_section["paragraphs"].append(line)

        if current_section["paragraphs"]:
            if "section_index" not in current_section:
                section_counter += 1
                current_section["section_index"] = section_counter
            sections.append(current_section)

        return sections
    
    def _extract_text_from_html(self, url: str) -> List[Dict[str, Any]]:
        sections =[]
        section_counter = 0
        try:
            loader = WebBaseLoader([url]) # WebBaseLoader expects a list of URLs
            documents = loader.load()
            for i,doc in enumerate(documents):
                text = doc.page_content.strip()
                if text:
                    title = doc.metadata.get("title", f"URL Content {i+1}")
                    current_section = {"title": title, "page_number": 1, "paragraphs": [text], "section_index": section_counter + 1}
                    sections.append(current_section)
                    section_counter += 1
        except Exception as e:
            print(f"Error in loading the html file {url}: {e}")
        return sections
    
    def extract_doc_id_from_url(self, url: str) -> str:
        """
        Extracts the Google Docs document ID from a full URL.
        """
        match = re.search(r"/d/([a-zA-Z0-9-_]+)", url)
        if match:
            return match.group(1)
        else:
            raise ValueError("Invalid Google Docs URL format.")


    def _extract_text_from_gdoc_url(self, url: str) -> List[Dict[str, Any]]:
        """
        Extracts text from a public Google Docs URL (exported as plain text).
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0

        doc_id = self.extract_doc_id_from_url(url)
        export_url = f"https://docs.google.com/document/d/{doc_id}/export?format=txt"
        response = requests.get(export_url)

        if response.status_code != 200:
            raise Exception(f"Failed to fetch Google Doc content: {response.status_code}")

        lines = [line.strip() for line in response.text.splitlines() if line.strip()]

        current_section = {
            "title": "page_1_untitled",
            "page_number": 1,
            "paragraphs": []
        }

        for line in lines:
            if self._is_heading(line):
                if current_section["paragraphs"]:
                    sections.append(current_section)
                section_counter += 1
                current_section = {
                    "title": line,
                    "page_number": 1,
                    "paragraphs": [],
                    "section_index": section_counter
                }
            else:
                current_section["paragraphs"].append(line)

        if current_section["paragraphs"]:
            if "section_index" not in current_section:
                section_counter += 1
                current_section["section_index"] = section_counter
            sections.append(current_section)

        return sections
    
    def _extract_text_from_gsheet_url(self, url: str) -> List[Dict[str, Any]]:
        """
        Extracting the data from google sheet (using url importing it as csv)
        RETURN a list of dict which contain data , id and more.
        """
        sections = []
        section_counter = 0

        # Convert the Google Sheets URL to CSV export format
        if "/edit#gid=" in url:
            base_url, gid_part = url.split("/edit#gid=")
            sheet_id = base_url.split("/d/")[1]
            gid = gid_part
        else:
            raise ValueError("Invalid Google Sheets URL format.")

        export_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv&gid={gid}"
        response = requests.get(export_url)

        if response.status_code != 200:
            raise Exception(f"Failed to fetch Google Sheet content: {response.status_code}")

        # Load into pandas DataFrame
        df = pd.read_csv(pd.compat.StringIO(response.text))

        current_section = {
            "title": "sheet_data",
            "page_number": 1,
            "paragraphs": []
        }

        for _, row in df.iterrows():
            # Create a readable paragraph from the row
            paragraph = "; ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val) and str(val).strip()])
            if paragraph:
                current_section["paragraphs"].append(paragraph)

        if current_section["paragraphs"]:
            section_counter += 1
            current_section["section_index"] = section_counter
            sections.append(current_section)

        return sections


    def _extract_text_from_image(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from an image using OCR.
        Treats the whole image as one section.
        """
        sections = []
        try:
            img = Image.open(file_path)
            with open(file_path, "rb") as f:
                        img_base64 = base64.b64encode(f.read()).decode("utf-8")
                        content = self._extract_text_from_vlm(encoded_image=img_base64)
                        if content:
                            sections.append({
                                "title": content.title,
                                "page_number": 1,
                                "paragraphs": [content.description],
                                "section_index": 1
                            })
                            return sections
                        
            text_from_ocr = pytesseract.image_to_string(img)
            if text_from_ocr.strip():
                # Split text into paragraphs
                paragraphs = [p.strip() for p in text_from_ocr.split('\n\n') if p.strip()]
                if not paragraphs:  # If no double newlines, try single newlines
                    paragraphs = [p.strip() for p in text_from_ocr.split('\n') if p.strip()]
                
                if paragraphs:
                    
                    sections.append({
                            "title": "ocr_image_section",
                            "page_number": 1,
                            "paragraphs": paragraphs,
                            "section_index": 1
                        })
            print(f"Extracted text using OCR from image: {os.path.basename(file_path)}")
        except pytesseract.TesseractNotFoundError:
            print("TesseractNotFoundError: Tesseract is not installed or not in your PATH.")
            print("Please install Tesseract OCR and ensure it's accessible.")
        except Exception as e:
            print(f"Error extracting text from image {os.path.basename(file_path)} using OCR: {e}")
        return sections

    def _perform_semantic_chunking(self, text: str) -> List[str]:
        """
        Perform semantic chunking on text using embeddings, fallback if needed.
        Here we use RecursiveCharacterTextSplitter for INSTRUCTOR-compatible chunking.
        """
        try:
            return self.semantic_splitter.split_text(text)
        except Exception as e:
            print(f"Chunking failed: {e}")
            return [text]  # fallback to whole text if even fallback splitter fails

    def process_document(self, file_path: str, source_doc_id: str, collection_name: Optional[str] = None) -> bool:
        """Process a document using fast rule-based chunking."""
        try:
            # Extract text based on file type
            file_extension = Path(file_path).suffix.lower()
            if file_extension == ".pdf":
                sections = self._extract_sections_from_pdf(file_path)
            elif file_extension == ".pptx":
                sections = self._extract_text_from_pptx(file_path)
            elif file_extension == ".csv":
                sections = self._extract_text_from_csv(file_path)
            elif file_extension in [".xlsx", ".xls"]:
                sections = self._extract_text_from_excel(file_path)
            elif file_extension == ".docx":
                sections = self._extract_text_from_docx(file_path)
            elif file_extension == ".json":
                sections = self._extract_text_from_json(file_path)
            elif file_extension == ".txt":
                sections = self._extract_text_from_txt(file_path)
            elif file_extension in [".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"]:
                sections = self._extract_text_from_image(file_path)
            else:
                print(f"Unsupported file type: {file_extension}")
                return False

            if not sections:
                print(f"No text extracted from {file_path}")
                return False

            # Process the text
            all_chunks, all_metadatas, all_ids = [], [], []
            global_chunk_counter = 1
            for sec in sections:
                section_text = "\n\n".join(sec["paragraphs"])
                chunks = self._perform_semantic_chunking(section_text)
                for idx, chunk in enumerate(chunks):
                    sec_idx = sec.get("section_index", sec["page_number"])
                    chunk_id = f"{source_doc_id}_sec{sec_idx}_chunk{global_chunk_counter}"
                    metadata = {
                        "source_doc_id": source_doc_id,
                        "file_name": Path(file_path).name,
                        "section_title": sec["title"],
                        "page_number": sec["page_number"],
                        "section_index": sec_idx,
                        "chunk_index": global_chunk_counter,
                        "paragraph_number_in_page": idx + 1
                    }
                    all_chunks.append(chunk)
                    all_metadatas.append(metadata)
                    all_ids.append(chunk_id)
                    global_chunk_counter += 1

            # Add documents to vector store
            success = self.vector_store_service.add_documents(
                chunks=all_chunks,
                metadatas=all_metadatas,
                doc_ids=all_ids,
                collection_name=collection_name
            )

            return success
        except Exception as e:
            print(f"Error processing document {file_path}: {e}")
            return False

    def process_url(self, url: str, source_doc_id: str, collection_name: Optional[str] = None) -> bool:
        """Process a document from a URL using fast rule-based chunking."""
        try:
            # Check if the URL is a Google Doc and route accordingly
            if "docs.google.com/document/" in url:
                print(f"Detected Google Doc URL: {url}")
                sections = self._extract_text_from_gdoc_url(url=url)
            elif "docs.google.com/spredsheet/" in url:
                print(f"Detected Google Sheet URL: {url}")
                sections = self._extract_text_from_gsheet_url(url=url)
            # elif "pdf" in url:
            #     print(f"Detected the pdf URL: {url}")
            #     sections = self._exatract_sections_from_onlinepdf(url=url)
            else:
                print(f"Detected standard web URL: {url}")
                sections = self._extract_text_from_html(url)

            if not sections:
                print(f"No text extracted from {url}")
                return False

            # Process the text
            all_chunks, all_metadatas, all_ids = [], [], []
            global_chunk_counter = 1
            for sec in sections:
                section_text = "\n\n".join(sec["paragraphs"])
                chunks = self._perform_semantic_chunking(section_text)
                for idx, chunk in enumerate(chunks):
                    sec_idx = sec.get("section_index", sec["page_number"])
                    chunk_id = f"{source_doc_id}_sec{sec_idx}_chunk{global_chunk_counter}"
                    metadata = {
                        "source_doc_id": source_doc_id,
                        "file_name": url,  # Use the URL as the identifier
                        "section_title": sec["title"],
                        "page_number": sec["page_number"],
                        "section_index": sec_idx,
                        "chunk_index": global_chunk_counter,
                        "paragraph_number_in_page": idx + 1
                    }
                    all_chunks.append(chunk)
                    all_metadatas.append(metadata)
                    all_ids.append(chunk_id)
                    global_chunk_counter += 1

            # Add documents to vector store
            success = self.vector_store_service.add_documents(
                chunks=all_chunks,
                metadatas=all_metadatas,
                doc_ids=all_ids,
                collection_name=collection_name
            )

            return success
        except Exception as e:
            print(f"Error processing URL {url}: {e}")
            return False

    def extract_text_only(self, file_path: str) -> str:
        """
        Extracts and returns all text from a document as a single string.
        Uses the same per-file-type extraction logic as process_document.
        """
        file_extension = Path(file_path).suffix.lower()
        if file_extension == ".pdf":
            sections = self._extract_sections_from_pdf(file_path)
        elif file_extension == ".pptx":
            sections = self._extract_text_from_pptx(file_path)
        elif file_extension == ".csv":
            sections = self._extract_text_from_csv(file_path)
        elif file_extension in [".xlsx", ".xls"]:
            sections = self._extract_text_from_excel(file_path)
        elif file_extension == ".docx":
            sections = self._extract_text_from_docx(file_path)
        elif file_extension == ".json":
            sections = self._extract_text_from_json(file_path)
        elif file_extension == ".txt":
            sections = self._extract_text_from_txt(file_path)
        elif file_extension in [".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"]:
            sections = self._extract_text_from_image(file_path)
        else:
            return "Unsupported file type for text extraction."

        if not sections:
            return "No text extracted from document."

        # Concatenate all paragraphs from all sections
        all_text = []
        for sec in sections:
            paragraphs = sec.get("paragraphs", [])
            if isinstance(paragraphs, list):
                all_text.extend(paragraphs)
            elif isinstance(paragraphs, str):
                all_text.append(paragraphs)
        return "\n\n".join(all_text)


def process_document_background(
    file_path: str, 
    source_doc_id: str, 
    doc_parser_svc_instance: DocParserFastService,
    serial_no: Optional[int] = None, 
    total_count: Optional[int] = None,
    collection_name: Optional[str] = None
):
    parser_name = doc_parser_svc_instance.__class__.__name__
    progress_log = f"Document {serial_no}/{total_count}" if serial_no and total_count else "Single document"
    
    print(f"Background task started for: {Path(file_path).name}, Source ID: {source_doc_id}, Parser: {parser_name}, ({progress_log})")
    try:
        success = doc_parser_svc_instance.process_document(
            file_path=file_path, 
            source_doc_id=source_doc_id,
            collection_name=collection_name
        )
        if success:
            print(f"Background processing completed successfully for {source_doc_id} ({Path(file_path).name}) using {parser_name}")
        else:
            print(f"Background processing (using {parser_name}) had issues or no chunks generated for {source_doc_id} ({Path(file_path).name})")
    except Exception as e:
        print(f"Error during background document processing for {source_doc_id} ({Path(file_path).name}) (using {parser_name}): {e}")


def process_url_background(
    url: str,
    source_doc_id: str,
    doc_parser_svc_instance: DocParserFastService,
    serial_no: Optional[int] = None,
    total_count: Optional[int] = None,
    collection_name: Optional[str] = None
):
    parser_name = doc_parser_svc_instance.__class__.__name__
    progress_log = f"URL {serial_no}/{total_count}" if serial_no and total_count else "Single URL"

    print(f"Background task started for: {url}, Source ID: {source_doc_id}, Parser: {parser_name}, ({progress_log})")
    try:
        success = doc_parser_svc_instance.process_url(
            url=url,
            source_doc_id=source_doc_id,
            collection_name=collection_name
        )
        if success:
            print(f"Background processing completed successfully for {source_doc_id} ({url}) using {parser_name}")
        else:
            print(f"Background processing (using {parser_name}) had issues or no chunks generated for {source_doc_id} ({url})")
    except Exception as e:
        print(f"Error during background document processing for {source_doc_id} ({url}) (using {parser_name}): {e}")


if __name__ == "__main__":
    process_document_background(file_path=r"C:\Users\hp\Downloads\Resume (1).docx")