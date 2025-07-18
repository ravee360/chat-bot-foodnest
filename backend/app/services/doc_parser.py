# backend/app/services/doc_parser.py

import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import io
import os
import uuid
from typing import List, Dict, Any, Optional, Tuple
import time # For potential retries or delays
from pathlib import Path

import requests # For calling OpenRouter API

#loading langchain_loaders to handle multiple files
from langchain_community.document_loaders import UnstructuredExcelLoader, CSVLoader, UnstructuredPowerPointLoader, Docx2txtLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from pptx import Presentation
import pandas as pd

# Local application imports
from backend.app.core.config import settings, PROJECT_ROOT_DIR
from backend.app.services.vstore_svc import VectorStoreService # Assuming singleton instance

# Optional: If Tesseract is not in your PATH, you might need to specify its location
# Example for Windows, adjust if necessary:
# TESSERACT_PATH = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
# if os.path.exists(TESSERACT_PATH):
#     pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH
# else:
#     print(f"Warning: Tesseract OCR path not found at {TESSERACT_PATH}. Ensure Tesseract is in your system PATH.")


class DocParserService:
    """
    Service to parse documents (PDFs, images), extract text, 
    perform LLM-based semantic chunking, and add chunks to the vector store.
    Now supports structure-aware parsing and multiple file types.
    """
    def __init__(self, vector_store_service: VectorStoreService, use_llm_chunking: bool = True):
        self.vector_store_service = vector_store_service
        self.use_llm_chunking = use_llm_chunking
        
        # Initialize embeddings for semantic chunking (fallback to LLM)
        try:
            self.embeddings = HuggingFaceEmbeddings(
                model_name="sentence-transformers/all-MiniLM-L6-v2",
                model_kwargs={'device': 'cpu'},
                encode_kwargs={"normalize_embeddings": True}
            )
            self.semantic_splitter = RecursiveCharacterTextSplitter(
                chunk_size=800,
                chunk_overlap=200,
                length_function=len,
                is_separator_regex=False,
            )
        except Exception as e:
            print(f"Error initializing embedding model: {e}. Using fallback splitter.")
            self.embeddings = None
            self.semantic_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
                is_separator_regex=False,
                separators=["\n\n", "\n", ". ", "? ", "! ", ",", " ", ""]
            )
        
        print(f"DocParserService initialized ({'LLM' if use_llm_chunking else 'Semantic'} Chunking Mode).")

    def _is_heading(self, text: str) -> bool:
        """
        Detect if a line/paragraph is a heading based on heuristic:
        - Mostly uppercase or capitalized and short (< 8 words)
        """
        words = text.strip().split()
        if len(words) <= 8 and sum(1 for w in words if w.isupper()) >= len(words) * 0.6:
            return True
        return False

    def _extract_sections_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from PDF and groups paragraphs into sections by detecting headings.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        doc = fitz.open(file_path)
        section_counter = 0
        for page_index in range(len(doc)):
            page = doc.load_page(page_index)
            blocks = page.get_text("blocks", sort=True)
            current_section = {"title": f"page_{page_index+1}_untitled", "page_number": page_index+1, "paragraphs": []}
            for block in blocks:
                if block[6] != 0:
                    continue
                text = block[4].replace('\r', ' ').replace('\n', ' ').strip()
                if not text:
                    continue
                if self._is_heading(text):
                    # start new section
                    if current_section["paragraphs"]:
                        sections.append(current_section)
                    section_counter += 1
                    current_section = {"title": text, "page_number": page_index+1, "paragraphs": [], "section_index": section_counter}
                else:
                    current_section["paragraphs"].append(text)
            if current_section["paragraphs"]:
                # ensure section_index for untitled sections
                if "section_index" not in current_section:
                    section_counter += 1
                    current_section["section_index"] = section_counter
                sections.append(current_section)
        doc.close()
        return sections

    def _extract_text_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from each page of a PDF.
        Each item in the list represents a page, containing a list of paragraphs.
        """
        pages_content = []
        try:
            doc = fitz.open(file_path)
            print(f"PDF '{os.path.basename(file_path)}' has {doc.page_count} pages.")
            for page_num_human_readable in range(1, doc.page_count + 1):
                page = doc.load_page(page_num_human_readable - 1) # fitz uses 0-indexed pages
                blocks = page.get_text("blocks", sort=True) # Sort by y-coordinate
                
                page_paragraphs_data = []
                para_counter_on_page = 0
                for block in blocks:
                    if block[6] == 0: # Text block
                        paragraph_text = block[4].replace('\r', ' ').replace('\n', ' ').strip()
                        if paragraph_text:
                            para_counter_on_page += 1
                            page_paragraphs_data.append({
                                "paragraph_number_in_page": para_counter_on_page,
                                "text": paragraph_text
                            })
                
                if page_paragraphs_data:
                    pages_content.append({
                        "page_number": page_num_human_readable,
                        "paragraphs": page_paragraphs_data
                    })
            doc.close()
            print(f"Extracted {len(pages_content)} pages with paragraphs from PDF: {os.path.basename(file_path)}")
        except Exception as e:
            print(f"Error extracting text from PDF {os.path.basename(file_path)}: {e}")
        return pages_content
    
    def _extract_text_from_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract the text from each slide of pptx 
        Each item in the list represents a slide, containing a list of paragraphs.
        """
        sections = []
        try:
            prs = Presentation(file_path)
            section_counter = 0

            for i, slide in enumerate(prs.slides):
                paragraphs = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            paragraphs.append(text)

                current_section = {
                    "title": f"slide_{i+1}_untitled",
                    "page_number": i + 1,
                    "paragraphs": paragraphs
                }
                if current_section["paragraphs"]:
                    # ensure section_index for untitled sections
                    if "section_index" not in current_section:
                        section_counter += 1
                        current_section["section_index"] = section_counter
                    sections.append(current_section)

            print(f"Extracted {len(sections)} slides with paragraphs from PPT: {os.path.basename(file_path)}")
        except Exception as e:
            print(f"Error in extracting the text from pptx document {os.path.basename(file_path)}: {e}")
        return sections
    

    def _extract_text_from_csv(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract the text from CSV file
        RET list of item containing page number with their paragraphs
        """
        sections = []
        section_counter = 0
        try:
            loader = CSVLoader(file_path=file_path)
            documents = loader.load_and_split(text_splitter=self.semantic_splitter)
            for i in range(len(documents)):
                current_section = {"title": f"page_{i+1}_untitled", "page_number": i+1, "paragraphs": []}
                current_section["paragraphs"].append(documents[i].page_content)
                if current_section["paragraphs"]:
                    # ensure section_index for untitled sections
                    if "section_index" not in current_section:
                        section_counter += 1
                        current_section["section_index"] = section_counter
                    sections.append(current_section)
            print(f"Extracted {len(documents)} lines of data from csv file {os.path.basename(file_path)}")
        except Exception as e:
            print(f"Error in extracting the text from csv file {os.path.basename(file_path)}: {e}")
        return sections

    def _extract_text_from_excel(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from Excel file and split into sections.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0
        try:
            excel_file = pd.ExcelFile(file_path)
                
            for i, sheet_name in enumerate(excel_file.sheet_names):
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                    
                # Convert to string format
                text_content = df.to_string(index=False)
                    
                current_section = {
                        "title": f"sheet_{i+1}",
                        "page_number": i+1,
                        "paragraphs": [text_content],
                    }
                if current_section["paragraphs"]:
                    # ensure section_index for untitled sections
                    if "section_index" not in current_section:
                        section_counter += 1
                        current_section["section_index"] = section_counter
                    sections.append(current_section)
                
            print(f"Extracted {len(sections)} sheets from Excel file {os.path.basename(file_path)}")
        except Exception as e:
            print(f"Error in extracting the text from Excel file {os.path.basename(file_path)}: {e}")
        return sections
    
    def _extract_text_from_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from DOCX file and split into sections.
        Returns list of sections with 'title', 'page_number', and 'paragraphs'.
        """
        sections = []
        section_counter = 0
        try:
            loader = Docx2txtLoader(file_path=file_path)
            documents = loader.load()
            
            for i, doc in enumerate(documents):
                text = doc.page_content.strip()
                if not text:
                    continue
                current_section = {"title": f"page_{i+1}_untitled", "page_number": i+1, "paragraphs": [text]}
                section_counter += 1
                current_section["section_index"] = section_counter
                sections.append(current_section)

            print(f"Extracted {len(sections)} sections from DOCX file {os.path.basename(file_path)}")
        except Exception as e:
            print(f"Error in extracting the text from DOCX file {os.path.basename(file_path)}: {e}")
        return sections

    def _extract_text_from_image_ocr(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extracts text from an image using OCR.
        Treats the whole image as one page. Attempts to split by double newlines as paragraphs.
        """
        pages_content = []
        try:
            img = Image.open(file_path)
            text_from_ocr = pytesseract.image_to_string(img)
            
            if text_from_ocr.strip():
                raw_paragraphs = [p.strip() for p in text_from_ocr.split('\n\n') if p.strip()]
                if not raw_paragraphs: # If no double newlines, try single newlines
                    raw_paragraphs = [p.strip() for p in text_from_ocr.split('\n') if p.strip()]

                page_paragraphs_data = []
                for i, para_text in enumerate(raw_paragraphs):
                    page_paragraphs_data.append({
                        "paragraph_number_in_page": i + 1,
                        "text": para_text
                    })

                if page_paragraphs_data:
                    pages_content.append({
                        "page_number": 1, # OCR'd image is considered a single page
                        "paragraphs": page_paragraphs_data
                    })
            print(f"Extracted text using OCR from image: {os.path.basename(file_path)}")
        except pytesseract.TesseractNotFoundError:
            print("TesseractNotFoundError: Tesseract is not installed or not in your PATH.")
            print("Please install Tesseract OCR and ensure it's accessible.")
        except Exception as e:
            print(f"Error extracting text from image {os.path.basename(file_path)} using OCR: {e}")
        return pages_content

    def _perform_semantic_chunking(self, text: str) -> List[str]:
        """
        Perform semantic chunking on text using embeddings, fallback if needed.
        Here we use RecursiveCharacterTextSplitter for INSTRUCTOR-compatible chunking.
        """
        try:
            return self.semantic_splitter.split_text(text)
        except Exception as e:
            print(f"Chunking failed: {e}")
            return [text]  # fallback to whole text if even fallback splitter fails

    def _chunk_paragraph_semantically_llm(self, paragraph_text: str, source_info: str) -> List[str]:
        """
        Chunks a single paragraph semantically using an LLM via OpenRouter.
        Returns a list of text chunks. If it fails for any reason, returns an empty list and logs a message.
        """
        if not settings.OPENROUTER_API_KEY:
            print(f"OpenRouter API key not configured. LLM semantic chunking skipped for: {source_info}")
            return []

        headers = {
            "Authorization": f"Bearer {settings.OPENROUTER_API_KEY}",
            "Content-Type": "application/json",
            "HTTP-Referer": settings.PROJECT_NAME, 
            "X-Title": settings.PROJECT_NAME 
        }
        
        prompt = f"""You are an expert text processor. Your task is to break down the following text from '{source_info}' into semantically coherent and self-contained chunks. Each chunk should ideally be between 2 to 5 sentences long and focus on a distinct idea or topic within the paragraph. Ensure that the chunks are grammatically correct and retain the original meaning and important entities. Do not summarize or add any information not present in the original text. Output EACH chunk on a new line. Do NOT add any extra formatting, numbering, or commentary before or after the chunks. Just provide the raw text chunks, each on its own line. If the input text is too short or cannot be meaningfully chunked according to these rules, return the original text as a single chunk.

Text to process:
\"\"\"
{paragraph_text}
\"\"\"

Semantically coherent chunks (each on a new line):"""

        payload = {
            "model": settings.DEFAULT_LLM_MODEL, 
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.2, 
            "max_tokens": 2000 
        }
        
        print(f"Attempting LLM semantic chunking for paragraph from {source_info} using {settings.DEFAULT_LLM_MODEL}...")
        try:
            response = requests.post(
                f"{settings.OPENROUTER_API_BASE}/chat/completions", 
                headers=headers, 
                json=payload, 
                timeout=90
            )
            response.raise_for_status() 
            
            response_json = response.json()
            if response_json.get("choices") and response_json["choices"][0].get("message"):
                content = response_json["choices"][0]["message"].get("content", "")
                llm_chunks = [chunk.strip() for chunk in content.split('\n') if chunk.strip()]
                
                if llm_chunks:
                    print(f"LLM generated {len(llm_chunks)} chunks for paragraph from {source_info}.")
                    return llm_chunks
                else:
                    print(f"LLM returned no usable chunks for {source_info}. Content: '{content}'. This paragraph will not be chunked.")
                    return [] 
            else:
                error_details = response_json.get("error", "Unknown error structure")
                print(f"LLM response malformed for {source_info}. Details: {error_details}. This paragraph will not be chunked.")
                return []
        except requests.exceptions.RequestException as e:
            print(f"API request failed for LLM semantic chunking ({source_info}): {e}. This paragraph will not be chunked.")
            return []
        except Exception as e:
            print(f"Unexpected error during LLM semantic chunking ({source_info}): {e}. This paragraph will not be chunked.")
            return []

    def process_document(self, file_path: str, source_doc_id: str, collection_name: Optional[str] = None) -> bool:
        """
        Processes a single document: extracts text, performs chunking (LLM or semantic), 
        and adds resulting chunks to the vector store.
        """
        print(f"\nProcessing document: {os.path.basename(file_path)} (Source ID: {source_doc_id}, Mode: {'LLM' if self.use_llm_chunking else 'Semantic'} Chunking)")
        file_extension = Path(file_path).suffix.lower()
        
        # Extract text based on file type
        extracted_content: List[Dict[str, Any]] = []
        if file_extension == ".pdf":
            if self.use_llm_chunking:
                extracted_content = self._extract_text_from_pdf(file_path)
            else:
                extracted_content = self._extract_sections_from_pdf(file_path)
        elif file_extension == ".pptx":
            extracted_content = self._extract_text_from_pptx(file_path)
        elif file_extension == ".csv":
            extracted_content = self._extract_text_from_csv(file_path)
        elif file_extension == ".xlsx":
            extracted_content = self._extract_text_from_excel(file_path)
        elif file_extension == ".docx":
            extracted_content = self._extract_text_from_docx(file_path)
        elif file_extension in [".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"]:
            extracted_content = self._extract_text_from_image_ocr(file_path)
        else:
            print(f"Unsupported file type: {file_extension} for {file_path}")
            return False

        if not extracted_content:
            print(f"No text could be extracted from {file_path}.")
            return False

        all_final_chunks_texts: List[str] = []
        all_final_chunks_metadatas: List[Dict[str, Any]] = []
        all_final_chunks_ids: List[str] = []
        
        total_paragraphs_processed = 0
        total_chunks_generated = 0

        if self.use_llm_chunking:
            # LLM-based chunking (original method)
            for page_content in extracted_content:
                page_num = page_content["page_number"]
                for paragraph_data in page_content.get("paragraphs", []):
                    total_paragraphs_processed += 1
                    para_num = paragraph_data["paragraph_number_in_page"]
                    para_text = paragraph_data["text"]
                    source_info_for_llm = f"Doc: {source_doc_id}, Page: {page_num}, Para: {para_num}"

                    current_paragraph_chunks: List[str] = []
                    current_paragraph_chunks = self._chunk_paragraph_semantically_llm(para_text, source_info_for_llm)
                    
                    if not current_paragraph_chunks:
                         print(f"LLM semantic chunking yielded no chunks for {source_info_for_llm} (or API key was missing/call failed). This paragraph is skipped.")
                    
                    for chunk_seq, chunk_text in enumerate(current_paragraph_chunks):
                        total_chunks_generated +=1
                        # Create a unique ID for each chunk
                        chunk_id = f"{source_doc_id}_p{page_num}_pr{para_num}_c{chunk_seq+1}"
                        metadata = {
                            "source_doc_id": source_doc_id,
                            "file_name": os.path.basename(file_path),
                            "page_number": page_num,
                            "paragraph_number_in_page": para_num, # Paragraph from which this chunk originated
                            "chunk_sequence_in_paragraph": chunk_seq + 1, # Sequence of this chunk within the original paragraph
                            "original_paragraph_text_preview": para_text[:150] + "..." # For context during review
                        }
                        all_final_chunks_texts.append(chunk_text)
                        all_final_chunks_metadatas.append(metadata)
                        all_final_chunks_ids.append(chunk_id)
        else:
            # Semantic chunking (new method)
            global_chunk_counter = 1
            for section in extracted_content:
                section_text = "\n\n".join(section["paragraphs"])
                chunks = self._perform_semantic_chunking(section_text)
                for idx, chunk in enumerate(chunks):
                    sec_idx = section.get("section_index", section["page_number"])
                    chunk_id = f"{source_doc_id}_sec{sec_idx}_chunk{global_chunk_counter}"
                    metadata = {
                        "source_doc_id": source_doc_id,
                        "file_name": Path(file_path).name,
                        "section_title": section["title"],
                        "page_number": section["page_number"],
                        "section_index": sec_idx,
                        "chunk_index": global_chunk_counter,
                        "paragraph_number_in_page": idx + 1
                    }
                    all_final_chunks_texts.append(chunk)
                    all_final_chunks_metadatas.append(metadata)
                    all_final_chunks_ids.append(chunk_id)
                    global_chunk_counter += 1
                    total_chunks_generated += 1
        
        print(f"\nFinished processing {os.path.basename(file_path)}.")
        print(f"Total paragraphs analyzed: {total_paragraphs_processed}")
        print(f"Total chunks generated: {total_chunks_generated}")

        if all_final_chunks_texts:
            print(f"Attempting to add {len(all_final_chunks_texts)} chunks to vector store...")
            success = self.vector_store_service.add_documents(
                chunks=all_final_chunks_texts,
                metadatas=all_final_chunks_metadatas,
                doc_ids=all_final_chunks_ids,
                collection_name=collection_name
            )
            if success:
                print(f"Successfully added chunks from {os.path.basename(file_path)} to vector store.")
                return True
            else:
                print(f"Failed to add chunks from {os.path.basename(file_path)} to vector store.")
                return False
        else:
            print(f"No chunks were produced from {os.path.basename(file_path)} to add to vector store.")
            return False

def process_document_background(
    file_path: str, 
    source_doc_id: str, 
    doc_parser_svc_instance: DocParserService,
    serial_no: Optional[int] = None, 
    total_count: Optional[int] = None,
    collection_name: Optional[str] = None
):
    """
    Background processing function for documents.
    """
    parser_name = doc_parser_svc_instance.__class__.__name__
    progress_log = f"Document {serial_no}/{total_count}" if serial_no and total_count else "Single document"
    
    print(f"Background task started for: {Path(file_path).name}, Source ID: {source_doc_id}, Parser: {parser_name}, ({progress_log})")
    try:
        success = doc_parser_svc_instance.process_document(
            file_path=file_path, 
            source_doc_id=source_doc_id,
            collection_name=collection_name
        )
        if success:
            print(f"Background processing completed successfully for {source_doc_id} ({Path(file_path).name}) using {parser_name}")
        else:
            print(f"Background processing (using {parser_name}) had issues or no chunks generated for {source_doc_id} ({Path(file_path).name})")
    except Exception as e:
        print(f"Error during background document processing for {source_doc_id} ({Path(file_path).name}) (using {parser_name}): {e}")

# --- Test Block ---
if __name__ == "__main__":
    print("--- Testing DocParserService (Enhanced with Multiple Chunking Modes) ---")

    print(f"OpenRouter API Key available for test: {'Yes' if settings.OPENROUTER_API_KEY else 'No'}")
    if not settings.OPENROUTER_API_KEY:
        print("WARNING: OpenRouter API Key is NOT configured. LLM chunking tests will be skipped or show no chunks.")
    print(f"Using LLM Model for test: {settings.DEFAULT_LLM_MODEL}")

    vstore_service_instance = VectorStoreService() 

    if not vstore_service_instance._langchain_chroma_instance:
        print("CRITICAL: VectorStoreService did not initialize correctly. Aborting test.")
    else:
        # Test both chunking modes
        for use_llm in [True, False]:
            chunking_mode = "LLM" if use_llm else "Semantic"
            print(f"\n--- Testing {chunking_mode} Chunking Mode ---")
            
            parser_service = DocParserService(
                vector_store_service=vstore_service_instance,
                use_llm_chunking=use_llm
            )

            test_data_dir = os.path.join(PROJECT_ROOT_DIR, "test_documents")
            os.makedirs(test_data_dir, exist_ok=True)

            # Create test files
            dummy_pdf_path = os.path.join(test_data_dir, f"dummy_{chunking_mode.lower()}.pdf")
            try:
                pdf_doc = fitz.open() 
                page = pdf_doc.new_page()
                page.insert_text((72, 72), "The first paragraph for chunking. It talks about AI and its impact on modern software development. This field is rapidly evolving.")
                page.insert_text((72, 144), "Another distinct idea: renewable energy sources are crucial for a sustainable future. Solar and wind power are leading examples.")
                page = pdf_doc.new_page()
                page.insert_text((72, 72), "Page two has a very short paragraph. Just this one sentence.")
                pdf_doc.save(dummy_pdf_path)
                pdf_doc.close()
                print(f"Created dummy PDF for {chunking_mode} test: {dummy_pdf_path}")
            except Exception as e:
                print(f"Could not create dummy PDF: {e}.")

            # Test document processing
            if os.path.exists(dummy_pdf_path):
                parser_service.process_document(
                    file_path=dummy_pdf_path,
                    source_doc_id=f"{chunking_mode.lower()}_pdf_test"
                )
            else:
                print(f"Skipping {chunking_mode} test: Dummy PDF not found at {dummy_pdf_path}")
        
        print("\n--- DocParserService Enhanced Test Complete ---")
 
